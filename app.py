import io
import json
import zipfile
from pathlib import Path
from datetime import datetime

import streamlit as st
from pptx import Presentation


def safe_name(name: str) -> str:
    return "".join(c if c.isalnum() or c in "-_." else "_" for c in name)


def extract_links_from_shape(shape):
    links = []
    if not hasattr(shape, "text_frame") or not shape.text_frame:
        return links
    for p in shape.text_frame.paragraphs:
        for run in p.runs:
            try:
                h = run.hyperlink
                if h and h.address:
                    links.append(h.address)
            except Exception:
                pass
    return links


def convert_pptx(uploaded_name: str, uploaded_bytes: bytes):
    prs = Presentation(io.BytesIO(uploaded_bytes))

    out = {
        "slides": [],
        "assets": [],
    }

    md_lines = [f"# {uploaded_name}", ""]

    # 이미지 추출용 zip 접근
    zf = zipfile.ZipFile(io.BytesIO(uploaded_bytes))

    image_index = 0

    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        links = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                t = shape.text.strip()
                if t:
                    texts.append(t)

            links.extend(extract_links_from_shape(shape))

            # 이미지 추출
            if shape.shape_type == 13:  # PICTURE
                try:
                    image = shape.image
                    ext = image.ext or "png"
                    image_index += 1
                    img_name = f"images/slide_{i:03d}_img_{image_index:03d}.{ext}"
                    out["assets"].append(
                        {
                            "type": "image",
                            "slide": i,
                            "file_path": img_name,
                            "content_type": image.content_type,
                            "size": len(image.blob),
                            "_blob": image.blob,
                        }
                    )
                except Exception:
                    pass

        notes_text = ""
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            notes_text = ""

        slide_title = ""
        try:
            if slide.shapes.title and slide.shapes.title.text:
                slide_title = slide.shapes.title.text.strip()
        except Exception:
            pass

        merged_text = "\n".join(t for t in texts if t)

        slide_rec = {
            "slide": i,
            "title": slide_title,
            "text": merged_text,
            "notes": notes_text,
            "links": sorted(set(links)),
            "source": uploaded_name,
        }
        out["slides"].append(slide_rec)

        md_lines.append(f"## Slide {i}: {slide_title if slide_title else '(untitled)'}")
        md_lines.append("")
        if merged_text:
            md_lines.append(merged_text)
            md_lines.append("")
        if slide_rec["links"]:
            md_lines.append("### Links")
            for u in slide_rec["links"]:
                md_lines.append(f"- {u}")
            md_lines.append("")
        if notes_text:
            md_lines.append("### Notes")
            md_lines.append(notes_text)
            md_lines.append("")

    # 첨부파일(embeddings/ole) raw 추출
    for name in zf.namelist():
        low = name.lower()
        if low.startswith("ppt/embeddings/") or low.startswith("ppt/oleobjects/"):
            raw = zf.read(name)
            asset_name = f"attachments/{safe_name(Path(name).name)}"
            out["assets"].append(
                {
                    "type": "attachment",
                    "slide": None,
                    "file_path": asset_name,
                    "source_path_in_pptx": name,
                    "size": len(raw),
                    "_blob": raw,
                }
            )

    manifest = {
        "source_file": uploaded_name,
        "generated_at": datetime.now().isoformat(),
        "slide_count": len(out["slides"]),
        "asset_count": len(out["assets"]),
        "schema": {
            "slides_jsonl": "slide,title,text,notes,links,source",
            "assets_jsonl": "type,slide,file_path,size,...",
        },
    }

    md_text = "\n".join(md_lines).strip() + "\n"

    # ZIP 결과물 생성
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as outzip:
        outzip.writestr("document.md", md_text)

        slides_jsonl = "\n".join(json.dumps(s, ensure_ascii=False) for s in out["slides"]) + "\n"
        outzip.writestr("slides.jsonl", slides_jsonl)

        assets_for_json = []
        for a in out["assets"]:
            copy = {k: v for k, v in a.items() if k != "_blob"}
            assets_for_json.append(copy)
        assets_jsonl = "\n".join(json.dumps(a, ensure_ascii=False) for a in assets_for_json) + ("\n" if assets_for_json else "")
        outzip.writestr("assets.jsonl", assets_jsonl)

        outzip.writestr("manifest.json", json.dumps(manifest, ensure_ascii=False, indent=2))

        for a in out["assets"]:
            blob = a.get("_blob")
            path = a.get("file_path")
            if blob and path:
                outzip.writestr(path, blob)

    buf.seek(0)
    return buf, manifest


st.set_page_config(page_title="PPTX 문서변환 에디터", layout="wide")
st.title("📄 PPTX → Markdown/DB 변환 에디터")
st.caption("링크/첨부/이미지를 포함해 AI 분석용 데이터셋으로 변환")

uploaded = st.file_uploader("PPTX 업로드", type=["pptx"])

if uploaded:
    st.success(f"업로드됨: {uploaded.name} ({uploaded.size:,} bytes)")

if st.button("변환 실행", type="primary", disabled=(uploaded is None)):
    with st.spinner("변환 중..."):
        zip_buf, manifest = convert_pptx(uploaded.name, uploaded.read())

    st.subheader("변환 결과")
    col1, col2 = st.columns(2)
    col1.metric("슬라이드 수", manifest["slide_count"])
    col2.metric("추출 자산 수", manifest["asset_count"])

    st.download_button(
        "결과 ZIP 다운로드",
        data=zip_buf,
        file_name=f"{Path(uploaded.name).stem}_converted_dataset.zip",
        mime="application/zip",
    )

    st.info("ZIP 안에 document.md / slides.jsonl / assets.jsonl / images/ / attachments/ 포함")
